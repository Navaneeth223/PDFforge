"""
converter.py — Office/HTML/OCR conversions.
All heavy operations run as subprocesses or in threads.
"""
import os
import asyncio
import io
import uuid
import zipfile
from typing import Optional, List
from PIL import Image

from services.storage import get_output_path, get_session_dir


# ─── LibreOffice: Office → PDF ─────────────────────────────────────────────────

async def office_to_pdf(session_id: str, input_path: str) -> str:
    """Convert Office document to PDF using LibreOffice headless (async subprocess) or pure-python fallback."""
    session_dir = get_session_dir(session_id)
    final_path = get_output_path(session_id, ".pdf")
    
    # Try LibreOffice first
    try:
        proc = await asyncio.create_subprocess_exec(
            "libreoffice",
            "--headless",
            "--convert-to", "pdf",
            "--outdir", session_dir,
            input_path,
            stdout=asyncio.subprocess.PIPE,
            stderr=asyncio.subprocess.PIPE,
        )
        try:
            stdout, stderr = await asyncio.wait_for(proc.communicate(), timeout=45)
            if proc.returncode == 0:
                stem = os.path.splitext(os.path.basename(input_path))[0]
                converted_path = os.path.join(session_dir, f"{stem}.pdf")
                if os.path.exists(converted_path):
                    os.rename(converted_path, final_path)
                    return final_path
        except Exception:
            try:
                proc.kill()
            except Exception:
                pass
    except Exception:
        pass  # Fallback to pure-Python

    # Fallback Mechanism
    ext = os.path.splitext(input_path)[1].lower()
    if ext in (".docx", ".doc"):
        import mammoth
        from xhtml2pdf import pisa
        with open(input_path, "rb") as docx_file:
            result = mammoth.convert_to_html(docx_file)
            html_content = result.value
        html_content = f"<html><head><meta charset='utf-8'><style>body {{ font-family: Helvetica, Arial, sans-serif; padding: 20px; }} p {{ margin-bottom: 12px; line-height: 1.4; }}</style></head><body>{html_content}</body></html>"
        with open(final_path, "wb") as result_file:
            pisa.CreatePDF(html_content, dest=result_file)
        return final_path
    
    elif ext in (".xlsx", ".xls", ".csv"):
        from xhtml2pdf import pisa
        if ext == ".csv":
            import csv
            html = ["<html><head><meta charset='utf-8'><style>body { font-family: Helvetica, Arial, sans-serif; } table { border-collapse: collapse; width: 100%; } td, th { border: 1px solid #ddd; padding: 6px; font-size: 10px; }</style></head><body><table>"]
            with open(input_path, mode='r', encoding='utf-8', errors='ignore') as f:
                reader = csv.reader(f)
                for row in reader:
                    html.append("<tr>")
                    for cell in row:
                        html.append(f"<td>{cell}</td>")
                    html.append("</tr>")
            html.append("</table></body></html>")
            html_content = "".join(html)
        else:
            from openpyxl import load_workbook
            wb = load_workbook(input_path, data_only=True)
            html = ["<html><head><meta charset='utf-8'><style>body { font-family: Helvetica, Arial, sans-serif; } table { border-collapse: collapse; width: 100%; margin-bottom: 20px; } th, td { border: 1px solid #ddd; padding: 8px; font-size: 10px; } th { background-color: #f2f2f2; text-align: left; }</style></head><body>"]
            for sheetname in wb.sheetnames:
                ws = wb[sheetname]
                html.append(f"<h2>{sheetname}</h2>")
                html.append("<table>")
                for row in ws.iter_rows(values_only=True):
                    if all(v is None for v in row):
                        continue
                    html.append("<tr>")
                    for cell in row:
                        val = "" if cell is None else str(cell)
                        html.append(f"<td>{val}</td>")
                    html.append("</tr>")
                html.append("</table>")
            html.append("</body></html>")
            html_content = "".join(html)
        with open(final_path, "wb") as f:
            pisa.CreatePDF(html_content, dest=f)
        return final_path

    elif ext in (".pptx", ".ppt"):
        from pptx import Presentation
        from xhtml2pdf import pisa
        prs = Presentation(input_path)
        html = ["<html><head><meta charset='utf-8'><style>body { font-family: Helvetica, Arial, sans-serif; } .slide { page-break-after: always; border: 1px solid #ccc; padding: 25px; margin-bottom: 20px; } h2 { color: #4f46e5; border-bottom: 1px solid #eee; padding-bottom: 8px; } p { font-size: 12px; line-height: 1.4; color: #374151; }</style></head><body>"]
        for i, slide in enumerate(prs.slides):
            html.append(f"<div class='slide'>")
            html.append(f"<h2>Slide {i + 1}</h2>")
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    html.append(f"<p>{shape.text.strip()}</p>")
            html.append("</div>")
        html.append("</body></html>")
        html_content = "".join(html)
        with open(final_path, "wb") as f:
            pisa.CreatePDF(html_content, dest=f)
        return final_path
        
    raise RuntimeError(f"Conversion failed: LibreOffice not found and format '{ext}' has no fallback.")


# ─── WeasyPrint: HTML → PDF ────────────────────────────────────────────────────

def html_to_pdf(session_id: str,
                html_content: Optional[str],
                url: Optional[str]) -> str:
    out_path = get_output_path(session_id, ".pdf")
    try:
        from weasyprint import HTML
        if url:
            HTML(url=url).write_pdf(out_path)
        elif html_content:
            HTML(string=html_content).write_pdf(out_path)
        else:
            raise ValueError("Either html_content or url must be provided.")
    except Exception:
        # Fallback to xhtml2pdf
        from xhtml2pdf import pisa
        if url:
            import requests
            resp = requests.get(url, timeout=30)
            content = resp.text
        elif html_content:
            content = html_content
        else:
            raise ValueError("Either html_content or url must be provided.")
        
        with open(out_path, "wb") as f:
            pisa.CreatePDF(content, dest=f)
    return out_path



# ─── PDF → Word (python-docx) ──────────────────────────────────────────────────

def pdf_to_word(session_id: str, input_path: str) -> str:
    import fitz
    from docx import Document
    from docx.shared import Pt

    doc = fitz.open(input_path)
    word = Document()
    for i, page in enumerate(doc):
        if i > 0:
            word.add_page_break()
        word.add_heading(f"Page {i + 1}", level=2)
        blocks = page.get_text("blocks")
        for block in sorted(blocks, key=lambda b: (b[1], b[0])):
            txt = block[4].strip()
            if txt:
                para = word.add_paragraph(txt)
                para.runs[0].font.size = Pt(10) if para.runs else None

    out_path = get_output_path(session_id, ".docx")
    word.save(out_path)
    doc.close()
    return out_path


# ─── PDF → Excel (pdfplumber) ──────────────────────────────────────────────────

def pdf_to_excel(session_id: str, input_path: str) -> str:
    import pdfplumber
    from openpyxl import Workbook

    wb = Workbook()
    wb.remove(wb.active)  # Remove default sheet

    with pdfplumber.open(input_path) as pdf:
        for i, page in enumerate(pdf.pages):
            ws = wb.create_sheet(title=f"Page {i + 1}")
            tables = page.extract_tables()
            row_offset = 1
            if tables:
                for table in tables:
                    for row in table:
                        ws.append([cell or "" for cell in row])
                    row_offset += len(table) + 2
            else:
                # No tables — just dump raw text in column A
                text = page.extract_text() or ""
                for line in text.splitlines():
                    ws.append([line])

    out_path = get_output_path(session_id, ".xlsx")
    wb.save(out_path)
    return out_path


# ─── PDF → PowerPoint (python-pptx) ───────────────────────────────────────────

def pdf_to_ppt(session_id: str, input_path: str) -> str:
    import fitz
    from pptx import Presentation
    from pptx.util import Inches

    doc = fitz.open(input_path)
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]  # blank layout
    mat = fitz.Matrix(2, 2)  # 2x scaling for crisp slides

    for i, page in enumerate(doc):
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img_buf = io.BytesIO(pix.tobytes("png"))
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(img_buf, 0, 0,
                                  width=prs.slide_width,
                                  height=prs.slide_height)

    out_path = get_output_path(session_id, ".pptx")
    prs.save(out_path)
    doc.close()
    return out_path


# ─── OCR (pytesseract) ─────────────────────────────────────────────────────────

def ocr_pdf(session_id: str, input_path: str, language: str, dpi: int) -> str:
    """Render each PDF page to an image and run Tesseract OCR, producing a searchable PDF."""
    import fitz
    import pytesseract

    doc = fitz.open(input_path)
    mat = fitz.Matrix(dpi / 72, dpi / 72)
    out_doc = fitz.open()

    for page in doc:
        pix = page.get_pixmap(matrix=mat, alpha=False)
        img = Image.frombytes("RGB", (pix.width, pix.height), pix.samples)
        pdf_bytes = pytesseract.image_to_pdf_or_hocr(img, extension="pdf", lang=language)
        tmp = fitz.open("pdf", pdf_bytes)
        out_doc.insert_pdf(tmp)
        tmp.close()

    out_path = get_output_path(session_id, ".pdf")
    out_doc.save(out_path, garbage=4, deflate=True)
    out_doc.close()
    doc.close()
    return out_path
